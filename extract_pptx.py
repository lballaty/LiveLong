#!/usr/bin/env python3

from pptx import Presentation
import sys

def extract_pptx_content(pptx_file):
    """Extract text content from PowerPoint presentation including notes"""
    try:
        prs = Presentation(pptx_file)

        content = []
        content.append("# PowerPoint Content Extraction")
        content.append(f"File: {pptx_file}")
        content.append(f"Total slides: {len(prs.slides)}")
        content.append("")

        for i, slide in enumerate(prs.slides, 1):
            content.append(f"## Slide {i}")

            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                content.append("### Slide Content:")
                for text in slide_text:
                    content.append(f"- {text}")
            else:
                content.append("### Slide Content:")
                content.append("- [No text content]")

            # Extract notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = []
                for shape in notes_slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Skip the slide content that gets repeated in notes
                        text = shape.text.strip()
                        if text and text not in slide_text:
                            notes_text.append(text)

                if notes_text:
                    content.append("### Speaker Notes:")
                    for note in notes_text:
                        content.append(f"- {note}")
                else:
                    content.append("### Speaker Notes:")
                    content.append("- [No notes]")
            else:
                content.append("### Speaker Notes:")
                content.append("- [No notes slide]")

            content.append("")

        return "\n".join(content)

    except Exception as e:
        return f"Error extracting PowerPoint content: {e}"

if __name__ == "__main__":
    pptx_file = "japanese_longevity_exercise.pptx"
    content = extract_pptx_content(pptx_file)
    print(content)